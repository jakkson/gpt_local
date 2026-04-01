"""
Batch document loader - reads PDF, DOCX, TXT, CSV, XLSX, PPTX, HTML, EML, MD
and image files (via OCR) and converts them to plain text for indexing.
Skips Dropbox online-only / offline placeholder files.
"""

import csv
import email
import io
import json
import logging
import os
import signal
from pathlib import Path
from typing import Generator

import chardet

from config import SUPPORTED_EXTENSIONS

logger = logging.getLogger(__name__)

# macOS UF_DATALESS flag — set on cloud-only placeholder files (Dropbox, iCloud)
_UF_DATALESS = 0x00000040

# Max seconds to spend on a single file (OCR, parsing, etc.)
FILE_TIMEOUT_SECONDS = 60
# Max file size to attempt OCR on (50 MB)
MAX_OCR_FILE_SIZE = 50 * 1024 * 1024


class FileTimeoutError(Exception):
    pass


def _timeout_handler(signum, frame):
    raise FileTimeoutError("File processing timed out")


def is_file_local(file_path: Path) -> bool:
    """
    Check if a file is actually available on disk (not a Dropbox/iCloud
    online-only placeholder). Returns False for offline files.
    """
    try:
        st = file_path.stat()

        # macOS: check the dataless flag (cloud-only placeholder)
        if hasattr(st, "st_flags") and (st.st_flags & _UF_DATALESS):
            return False

        # Skip zero-byte files
        if st.st_size == 0:
            return False

        # Safety: try reading 1 byte to confirm the file is readable
        with open(file_path, "rb") as f:
            f.read(1)

        return True

    except OSError:
        return False


class Document:
    """Simple container for a parsed document."""

    def __init__(self, text: str, metadata: dict):
        self.text = text
        self.metadata = metadata

    def __repr__(self):
        name = self.metadata.get("filename", "unknown")
        return f"Document({name}, {len(self.text)} chars)"


def detect_encoding(file_path: Path) -> str:
    with open(file_path, "rb") as f:
        raw = f.read(10_000)
    result = chardet.detect(raw)
    return result.get("encoding", "utf-8") or "utf-8"


def load_txt(file_path: Path) -> str:
    enc = detect_encoding(file_path)
    return file_path.read_text(encoding=enc, errors="replace")


def load_image(file_path: Path) -> str:
    """OCR an image file using Tesseract, with size guard and timeout."""
    if file_path.stat().st_size > MAX_OCR_FILE_SIZE:
        logger.warning(f"Skipping oversized image: {file_path.name}")
        return ""

    import pytesseract
    from PIL import Image

    img = Image.open(file_path)
    if img.mode != "RGB":
        img = img.convert("RGB")

    width, height = img.size
    max_dim = 4000
    if width > max_dim or height > max_dim:
        ratio = min(max_dim / width, max_dim / height)
        img = img.resize((int(width * ratio), int(height * ratio)))

    text = pytesseract.image_to_string(img, lang="eng", timeout=30)
    return text.strip()


def load_pdf(file_path: Path) -> str:
    """Extract text from PDF; falls back to OCR for scanned/image-based pages."""
    from pypdf import PdfReader

    reader = PdfReader(str(file_path))
    pages = []
    ocr_pages = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(text)
        else:
            ocr_pages.append(i)

    if ocr_pages:
        try:
            import pytesseract
            from pdf2image import convert_from_path

            images = convert_from_path(str(file_path), dpi=300)
            for i in ocr_pages:
                if i < len(images):
                    ocr_text = pytesseract.image_to_string(
                        images[i], lang="eng", timeout=30
                    )
                    if ocr_text and ocr_text.strip():
                        pages.insert(i, f"[OCR Page {i + 1}]\n{ocr_text.strip()}")
            logger.info(f"OCR'd {len(ocr_pages)} scanned pages in {file_path.name}")
        except Exception as e:
            logger.warning(f"OCR fallback failed for {file_path.name}: {e}")

    return "\n\n".join(pages)


def load_docx(file_path: Path) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def load_xlsx(file_path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def load_csv(file_path: Path) -> str:
    enc = detect_encoding(file_path)
    with open(file_path, "r", encoding=enc, errors="replace") as f:
        reader = csv.reader(f)
        rows = []
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def load_pptx(file_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def load_html(file_path: Path) -> str:
    from bs4 import BeautifulSoup

    enc = detect_encoding(file_path)
    raw = file_path.read_text(encoding=enc, errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def load_eml(file_path: Path) -> str:
    raw = file_path.read_bytes()
    msg = email.message_from_bytes(raw)
    parts = []
    parts.append(f"From: {msg.get('From', '')}")
    parts.append(f"To: {msg.get('To', '')}")
    parts.append(f"Subject: {msg.get('Subject', '')}")
    parts.append(f"Date: {msg.get('Date', '')}")
    parts.append("")

    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            if ctype == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="replace"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            parts.append(payload.decode("utf-8", errors="replace"))

    return "\n".join(parts)


def load_json(file_path: Path) -> str:
    enc = detect_encoding(file_path)
    raw = file_path.read_text(encoding=enc, errors="replace")
    try:
        data = json.loads(raw)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return raw


def load_md(file_path: Path) -> str:
    return load_txt(file_path)


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".doc": load_docx,
    ".txt": load_txt,
    ".md": load_md,
    ".csv": load_csv,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".pptx": load_pptx,
    ".html": load_html,
    ".htm": load_html,
    ".eml": load_eml,
    ".json": load_json,
    ".xml": load_txt,
    ".rtf": load_txt,
    ".log": load_txt,
    # Image files (OCR)
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".tiff": load_image,
    ".tif": load_image,
    ".bmp": load_image,
    ".gif": load_image,
    ".webp": load_image,
    ".heic": load_image,
}


def load_single_file(file_path: Path) -> Document | None:
    """Load a single file and return a Document, or None on failure."""
    ext = file_path.suffix.lower()
    if ext not in LOADERS:
        logger.warning(f"Unsupported file type: {ext} ({file_path.name})")
        return None

    if not is_file_local(file_path):
        logger.info(f"Skipping offline file: {file_path.name}")
        return None

    old_handler = signal.signal(signal.SIGALRM, _timeout_handler)
    signal.alarm(FILE_TIMEOUT_SECONDS)

    try:
        loader = LOADERS[ext]
        text = loader(file_path)
        signal.alarm(0)

        if not text or not text.strip():
            logger.warning(f"Empty content: {file_path.name}")
            return None

        metadata = {
            "filename": file_path.name,
            "filepath": str(file_path.resolve()),
            "extension": ext,
            "size_bytes": file_path.stat().st_size,
        }
        return Document(text=text.strip(), metadata=metadata)

    except FileTimeoutError:
        logger.warning(f"Timed out after {FILE_TIMEOUT_SECONDS}s: {file_path.name}")
        return None
    except OSError as e:
        logger.info(f"Skipping unavailable file: {file_path.name} ({e})")
        return None
    except Exception as e:
        logger.error(f"Failed to load {file_path.name}: {e}")
        return None
    finally:
        signal.alarm(0)
        signal.signal(signal.SIGALRM, old_handler)


def scan_directory(directory: Path, recursive: bool = True) -> list[Path]:
    """Find all supported files in a directory, skipping offline/cloud-only files."""
    files = []
    skipped = 0
    pattern = "**/*" if recursive else "*"
    for path in directory.glob(pattern):
        if not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if path.name.startswith("."):
            continue
        if not is_file_local(path):
            skipped += 1
            continue
        files.append(path)
    if skipped:
        logger.info(f"Skipped {skipped} offline/cloud-only files in {directory}")
    return sorted(files)


def load_documents(
    paths: list[Path], recursive: bool = True
) -> Generator[Document, None, None]:
    """
    Batch load documents from a mix of files and directories.
    Yields Document objects as they're parsed.
    """
    for path in paths:
        if path.is_dir():
            files = scan_directory(path, recursive=recursive)
            logger.info(f"Found {len(files)} files in {path}")
            for f in files:
                doc = load_single_file(f)
                if doc:
                    yield doc
        elif path.is_file():
            doc = load_single_file(path)
            if doc:
                yield doc
        else:
            logger.warning(f"Path not found: {path}")
